from pptx import Presentation;
from pptx.util import Inches;
from pptx.enum.shapes import MSO_SHAPE;
from pptx.dml.color import ColorFormat, RGBColor;
from pptx.enum.text import PP_ALIGN;
import glob;
import PIL;
from PIL import Image;
from PIL.ExifTags import TAGS;
from io import BytesIO;
from datetime import datetime;
import unicodedata;

# Retrieve all img paths! list in list
IMG_PATH = "G:\ATD\Signs_and_Markings\MISC_PROJECTS\\"
img_paths = [glob.glob(IMG_PATH + "Photos_TechRoom_TV\TV Use\**\*.jpg")] # This one searches all subfolders

def split_mpo(filename):
    with open(filename, 'rb') as f:
        data = f.read()
        # Look for the hex string 0xFFD9FFD8FFE1:
        #   0xFFD9 represents the end of the first JPEG image
        #   0xFFD8FFE1 marks the start of the appended JPEG image
        idx = data.find(b'\xFF\xD8\xFF\xE1', 1)
        if idx > 0:
            file = Image.open(BytesIO(data[: idx]))
            newfilepath = filename.split("\\")
            newfilepath = IMG_PATH + r'\Photos_TechRoom_TV\MPO_Flat\\'+ str(newfilepath[-1])
            if file is not None:
                file.save(newfilepath)
                file.close()
                return newfilepath
        else:
            return filename


def exifHarvest(filename):
    # These variables must declared first
    image = Image.open(filename)
    metadata = image.getexif()
    xpComment = ""
    orientation = 1
    date = ""

    if metadata is not None:

        for tag_id in metadata:
            # get the tag name, instead of human unreadable tag id
            tag = TAGS.get(tag_id, tag_id)
            # print(tag) # tests what tags available
            if tag == "ImageDescription" or tag == "XPComment":
                data = metadata.get(tag_id)
                if isinstance(data, bytes):
                    data = data.decode('latin-1')  # .decode('utf-8') #latin-1
                # print(f"{tag:25}, {data}")
                xpComment = data
            if tag == "Orientation":
                orientation = metadata.get(tag_id)
            if tag == "DateTime":
                date = metadata.get(tag_id)
                date = datetime.strptime(date, '%Y:%m:%d %H:%M:%S').strftime("%B %d %Y")
        # imgSize = (metadata.get("ImageWidth"), metadata.get("ImageLength"))
        # imgComment = metadata.get(0x010e)
        # print(metadata)
        return {"size": image.size, "comment": xpComment, "orientation": orientation, "date": date}

def remove_control_characters(s):
    return "".join(ch for ch in s if unicodedata.category(ch)[0] != "C")

# Do this
error_paths = []
prs = Presentation()
width = Inches(13.33)
height = Inches(7.5)
prs.slide_width = width
prs.slide_height = height
blank_slide_layout = prs.slide_layouts[6]
for img_path in img_paths:
    for i in img_path:
        x = split_mpo(i)
        slide = prs.slides.add_slide(blank_slide_layout)

        # Make the slide a black color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        try:
            metadata = exifHarvest(x)
            top = Inches(0)

            shape = slide.shapes

            # Size photo
            imgratio = metadata["size"][0] / metadata["size"][1]
            localheight = height

            # if pic is rotated
            if (metadata["orientation"] == 6 or metadata["orientation"] == 8):
                localheight = localheight / imgratio
                top = (height - localheight) / 2
                # imgratio = 1/imgratio

            # move image up to make room for comment
            # localheight = localheight - Inches(0.5)

            # Add comment
            # space to left, height, width of text box, length of text box
            #title = shape.add_textbox(Inches(0.25), height - Inches(0.75), width - Inches(0.5), Inches(1))
            #tf = title.text_frame
            #title.word_wrap = False
            #tf.clear()
            #p = tf.paragraphs[0]
            #p.alignment = PP_ALIGN.CENTER
            #run = p.add_run()
            #run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            #comment = remove_control_characters(metadata["comment"])
            #date_comment = remove_control_characters(metadata["date"])
            #run.text = comment + " ({})".format(date_comment)
            # hlink = run.hyperlink
            # hlink.address = i.rsplit('\\',1)[0]

            # Add centered photo
            left = (width - localheight * imgratio) / 2
            pic = shape.add_picture(x, left, top, localheight * imgratio, localheight)

            # Rotate if needed
            if (metadata["orientation"] == 6):
                pic.rotation = 90
            if (metadata["orientation"] == 8):
                pic.rotation = 270

        except Exception as e:
            print(e)
            error_paths.append(i)
            pass
prs.save(r'G:\ATD\Signs_and_Markings\MISC_PROJECTS\Photos_TechRoom_TV\TV.pptx')
print("done")